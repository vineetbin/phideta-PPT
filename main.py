from phi.agent.python import PythonAgent
from phi.model.openai import OpenAIChat
import json
from pptx import Presentation
from pptx.util import Pt

class AgentInitializer:
    """Initializes the LLM agent with the required settings."""
    
    def __init__(self):
        self.agent = self._initialize_agent()

    def _initialize_agent(self):
        """Sets up the PythonAgent with predefined parameters."""
        return PythonAgent(
            model=OpenAIChat(id="gpt-4o"),
            markdown=True,
            pip_install=True,
            show_tool_calls=True,
            add_chat_history_to_messages=True,
            num_history_messages=6,
        )

    def get_agent(self):
        return self.agent

class FileManager:
    """Handles file loading and saving operations."""
    
    @staticmethod
    def load_text_file(file_path: str) -> str:
        with open(file_path, 'r') as file:
            return file.read()

    @staticmethod
    def load_json_file(file_path: str) -> dict:
        with open(file_path, 'r') as file:
            return json.load(file)

    @staticmethod
    def save_text_file(file_path: str, content: str):
        with open(file_path, 'w') as file:
            file.write(content)

class PromptProcessor:
    """Processes prompts with the LLM agent and saves responses."""
    
    def __init__(self, agent, prompts: list):
        self.agent = agent
        self.prompts = prompts

    def process_prompts(self):
        """Processes each prompt, generating and saving responses."""
        for idx, prompt_data in enumerate(self.prompts):
            prompt = prompt_data.get("prompt", "")
            output_filename = f"response_{idx + 1}.txt"
            response = self._generate_response(prompt)
            FileManager.save_text_file(output_filename, response)
            print(f"Saved response to {output_filename}")

    def _generate_response(self, prompt: str) -> str:
        """Generates a response from the agent based on the prompt."""
        response = self.agent.run(prompt)
        return response.get_content_as_string()

class PresentationCreator:
    """Creates a PowerPoint presentation from text responses."""
    
    def __init__(self, agent, num_slides: int = 4):
        self.agent = agent
        self.num_slides = num_slides
        self.presentation = Presentation()

    def generate_presentation(self, pptx_filename: str = "generated_presentation.pptx"):
        """Generates a PowerPoint presentation based on response files."""
        for idx in range(1, self.num_slides + 1):
            slide_content = FileManager.load_text_file(f"response_{idx}.txt")
            formatted_content = self._get_formatted_content(slide_content)
            self._create_slide(formatted_content)
        self.presentation.save(pptx_filename)
        print(f"PowerPoint presentation saved as {pptx_filename}")

    def _get_formatted_content(self, slide_content: str) -> str:
        """Formats slide content by requesting structured output from the agent."""
        llm_prompt = (
            f"Format the below given content into a well-structured PowerPoint slide. "
            f"Mind to use only one slide for one .txt file and keep the height and width "
            f"across all slides even. I will be importing it to a google slide, make sure "
            f"no bugs are there. Also give only the required answer and no extra lines, "
            f"I will directly be using your whole answer in the slide:\n\n{slide_content}"
        )
        response = self.agent.run(llm_prompt)
        return response.get_content_as_string()

    def _create_slide(self, formatted_content: str):
        """Creates a slide with title and body based on formatted content."""
        slide_layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title, body = formatted_content.split("\n", 1)
        title_placeholder.text = title.strip()

        text_frame = content_placeholder.text_frame
        text_frame.text = body.strip()

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.name = "Arial"
            paragraph.space_after = Pt(10)

class ChatSession:
    """Coordinates the entire chat session from background setup to prompt processing."""
    
    def __init__(self, background_file: str, prompt_file: str):
        self.background_file = background_file
        self.prompt_file = prompt_file
        self.agent_initializer = AgentInitializer()
        self.agent = self.agent_initializer.get_agent()
        self.prompts = self._load_prompts()

    def _load_prompts(self) -> list:
        """Loads prompt data from a JSON file."""
        return FileManager.load_json_file(self.prompt_file)

    def run(self):
        """Executes the chat session workflow: loads background, processes prompts, and creates presentation."""
        self._load_background()
        prompt_processor = PromptProcessor(self.agent, self.prompts)
        prompt_processor.process_prompts()
        
        presentation_creator = PresentationCreator(self.agent, num_slides=len(self.prompts))
        presentation_creator.generate_presentation()

    def _load_background(self):
        """Loads and sets the background context for the agent."""
        background_text = FileManager.load_text_file(self.background_file)
        self.agent.print_response(f"Background: {background_text}")
        print("System initialized with background context.")

if __name__ == "__main__":
    background_file = 'constants/background_info.txt'
    prompt_file = 'constants/prompt.json'

    chat_session = ChatSession(background_file, prompt_file)
    chat_session.run()
